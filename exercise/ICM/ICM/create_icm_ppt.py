# -*- coding: utf-8 -*-
"""ICM20948 STM32 학습/분석 내용을 PPT로 정리하는 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
    return slide

def add_section_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    p = tx.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return slide

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 1. 표지
add_title_slide(prs,
    "STM32 ICM20948 학습 및 코드 분석",
    "데이터시트 공부 · 드라이버 분석 · ICM_test1 프로젝트\nD:\\git_file\\STM32\\STM32-study\\exercise\\ICM\\ICM"
)

# 2. 목차
add_content_slide(prs, "목차", [
    "1. 프로젝트 개요 및 사용 자료",
    "2. ICM20948 데이터시트 공부 (note.md)",
    "3. 코드 분석 (code_Analysis, ICM20948_코드분석.md)",
    "4. ICM_test1 STM32 IDE 프로젝트",
    "5. 정리 및 활용"
])

# 3. 사용한 파일/프로젝트
add_content_slide(prs, "1. 프로젝트 개요 · 사용 자료", [
    "• note.md : ICM20948 데이터시트 공부 필기 (Bank, DLPF, LSB, FSR, PLL, dps, 레지스터 등)",
    "• code_Analysis (1).md : ICM20948 헤더/소스 코드 분석 (매크로, 구조체, 함수 설명)",
    "• ICM20948_코드분석.md (ICM_test1 내) : 주소·레지스터·함수별 정리, MPU9250 비교, Q&A",
    "• ICM_test1 : STM32CubeIDE 프로젝트 (STM32F411RE, I2C1, UART2, ICM20948 드라이버)",
    "※ MPU_STM32, README는 해당 경로에서 미확인 — 추가 시 슬라이드 보강 가능"
])

# 4. 데이터시트 공부 요약
add_content_slide(prs, "2. 데이터시트 공부 요약 (note.md)", [
    "• Bank : 레지스터를 묶은 '페이지'. 0x7F에 0x00/0x20/0x30 등으로 Bank 선택.",
    "• Bank0=데이터/전원, Bank2=자이로·가속도 설정, Bank3=I2C 마스터(마그 AK09916 접근).",
    "• DLPF : 디지털 저역통과 필터. 노이즈 감소, 대신 지연·대역폭 트레이드오프.",
    "• LSB : raw 값의 '한 칸' 크기. 데이터시트의 µT/LSB, LSB/g 등으로 물리량 변환.",
    "• FSR : 측정 범위(±2g, ±500dps 등). FSR에 따라 스케일(8192, 65.5 등) 결정.",
    "• PLL : 센서 내부 클럭 안정화. I2C SCL과 별개.",
    "• dps : 자이로 단위(초당 각도). 자이로는 '속도'만 주고, 각도는 적분으로 구함."
])

# 5. 레지스터/실무 팁
add_content_slide(prs, "2. 데이터시트 공부 요약 (계속)", [
    "• PWR_MGMT_2 : 가속도/자이로 축별(X,Y,Z) 전원 ON/OFF. 실무에선 보통 전체 ON/OFF.",
    "• AK09916 ST1/ST2 : 마그 데이터 준비·읽기 시작(ST1), 읽기 완료·오버플로(ST2).",
    "• I2C 주소 : AD0=GND → 0x68, AD0=VCC → 0x69. Init에서 0x68 먼저 시도 후 0x69.",
    "• 실무 실수 : Bank 안 바꾸고 레지스터 접근 금지. select_bank → 쓰기/읽기 → 필요 시 Bank0 복귀."
])

# 6. 코드 분석 요약
add_content_slide(prs, "3. 코드 분석 요약 (code_Analysis, ICM20948_코드분석)", [
    "• SelectBank / WriteReg / ReadRegs : 0x7F로 Bank 선택 후 HAL_I2C_Mem_Write/Read 사용.",
    "• ICM20948_Init : 0x68→0x69 시도, WHO_AM_I(0xEA) 확인, 리셋·PLL·축 ON, Bank2 설정(4g/500dps),",
    "  I2C 마스터 활성화, Bank3에서 AK09916 리셋·연속 모드 설정 후 Bank0 복귀.",
    "• ReadAccel / ReadGyro : Bank0에서 0x2D~0x32(가속도), 0x33~0x38(자이로) 6바이트 읽어 raw·물리량 변환.",
    "• ReadMag : Bank3 SLV0로 AK09916 ST1부터 8바이트 읽기 설정 → Bank0 0x3B에서 결과 읽어 raw·µT 계산.",
    "• ICM20948 vs MPU9250 : WHO_AM_I 주소/값, 데이터 주소, Bank 유무 등 차이 정리됨."
])

# 7. ICM_test1 프로젝트
add_content_slide(prs, "4. ICM_test1 STM32 IDE 프로젝트", [
    "• MCU : STM32F411RE. HAL, I2C1(400kHz), UART2(시리얼 출력).",
    "• Core : main.c에서 ICM20948_Init → 실패 시 LD2 깜빡이며 대기, 성공 시 루프에서 ReadAll 호출.",
    "• ReadAll로 가속도(raw, g), 자이로(raw, dps), 마그(raw, uT) 읽어 UART로 출력.",
    "• 드라이버 : Core/Inc·Src의 icm20948.h / icm20948.c (Bank 선택, Init, ReadAccel/Gyro/Mag, ReadAll).",
    "• 정리 : 데이터시트·코드 분석을 바탕으로 ICM20948 + AK09916 9축 데이터를 STM32에서 읽고 시리얼로 확인한 프로젝트."
])

# 8. 정리
add_content_slide(prs, "5. 정리 및 활용", [
    "• 데이터시트(note) : Bank, DLPF, LSB, FSR, PLL, dps, 레지스터 역할을 정리해 두어 코드 이해에 활용.",
    "• 코드 분석 문서 : 헤더/소스 매크로·함수별 설명, MPU9250 비교, Q&A로 실습 시 참고.",
    "• ICM_test1 : 실제 보드에서 I2C 연결 후 Init·ReadAll로 9축 값 확인 가능.",
    "• 추가로 MPU_STM32 필기·README가 있으면 같은 형식으로 슬라이드 추가하면 됨."
])

# 9. 마무리
add_title_slide(prs, "감사합니다", "질문 있으시면 편하게 말씀해 주세요.")

out_path = r"D:\git_file\STM32\STM32-study\exercise\ICM\ICM\ICM20948_학습_발표.pptx"
prs.save(out_path)
print("저장 완료:", out_path)
